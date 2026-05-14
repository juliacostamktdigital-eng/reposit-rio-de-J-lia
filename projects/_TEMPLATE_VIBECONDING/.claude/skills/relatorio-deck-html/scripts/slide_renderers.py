"""
slide_renderers.py — Biblioteca de funções para construir cada tipo de slide
do deck V4 Colli em python-pptx.

Cada função recebe:
    prs   : Presentation  — objeto principal
    data  : dict          — JSON completo (replace + raw)
    logo  : str|None      — caminho para logov4.png (já convertida de webp)

E retorna o slide recém-criado.

Paridade com deck-base.html: cada função replica layout, cores e tipografia
do CSS V4 Colli (preto + vermelho #dc2626).
"""

import html
import io
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

import design_tokens as T

# ---------------------------------------------------------------------------
# Helpers internos
# ---------------------------------------------------------------------------

def _blank_slide(prs):
    """Adiciona e retorna slide em branco (layout 6 = blank)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _set_slide_bg(slide, light=False):
    """Preenche o fundo do slide com cor sólida V4."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = T.C_LT_BG if light else T.C_BG_DEEP


def _add_text(slide, text, x, y, w, h, *,
              font_name=T.FONT_PRIMARY,
              font_size=T.PT_BODY,
              bold=False,
              italic=False,
              color=None,
              align=PP_ALIGN.LEFT,
              wrap=True,
              word_wrap=True,
              margin_l=Inches(0.05),
              margin_t=Inches(0.03),
              margin_r=Inches(0.05),
              margin_b=Inches(0.03),
              auto_fit=False):
    """Adiciona caixa de texto simples."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left   = margin_l
    tf.margin_top    = margin_t
    tf.margin_right  = margin_r
    tf.margin_bottom = margin_b

    if auto_fit:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _strip_html(text) if _has_html(text) else text

    font = run.font
    font.name = font_name
    font.size = font_size
    font.bold = bold
    font.italic = italic
    if color:
        font.color.rgb = color
    return txBox


def _add_shape_rect(slide, x, y, w, h, *, fill_color, line_color=None, line_width=Pt(0.5),
                    rounded=False):
    """Adiciona retângulo preenchido com cor sólida. rounded=True aplica cantos arredondados."""
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE = 1
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    if rounded:
        # Aplica roundrect via ajuste de avLst no XML (equivalente a border-radius: 10px)
        sp_pr = shape._element.find(qn('p:spPr'))
        if sp_pr is None:
            sp_pr = shape._element.spPr
        prstGeom = sp_pr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            # adj controla o raio: 16667 ≈ 10% (equivale a border-radius proporcional)
            gd = avLst.find(qn('a:gd'))
            if gd is None:
                gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', 'val 12000')
    return shape


def _add_accent_bar(slide, x, y, w=Inches(0.9), h=Inches(0.06)):
    """Barra accent-bar vermelha horizontal (como .accent-bar no CSS)."""
    bar = slide.shapes.add_shape(1, x, y, w, h)
    fill = bar.fill
    fill.solid()
    fill.fore_color.rgb = T.C_ACCENT
    bar.line.fill.background()
    return bar


def _add_slide_label(slide, label, x, y, w, *, light=False):
    """Adiciona slide-label: texto em maiúsculas com barra vermelha à esquerda."""
    bar_w = Inches(0.04)
    bar_h = Inches(0.18)
    bar = slide.shapes.add_shape(1, x, y + Inches(0.01), bar_w, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = T.C_ACCENT
    bar.line.fill.background()

    txt_color = T.C_ACCENT_BR if not light else T.C_ACCENT
    _add_text(slide, label.upper(),
              x + bar_w + Inches(0.06), y,
              w - bar_w - Inches(0.06), Inches(0.22),
              font_size=T.PT_LABEL,
              bold=True,
              color=txt_color,
              align=PP_ALIGN.LEFT)


def _add_logo(slide, logo_path, light=False):
    """Insere logo no canto inferior direito (equivalente a .slide-logo)."""
    if not logo_path or not Path(logo_path).exists():
        return
    logo_h = Inches(0.28)
    logo_w = Inches(1.5)
    x = T.SLIDE_W - logo_w - Inches(0.18)
    y = T.SLIDE_H - logo_h - Inches(0.14)
    try:
        slide.shapes.add_picture(logo_path, x, y, width=logo_w, height=logo_h)
    except Exception:
        pass


def _slot_box(slide, x, y, w, h, *, light=False, dash=True):
    """
    Caixa estilo .slot:
      - fundo levemente tintado
      - borda accent (dashed via linha sólida fina — OOXML não suporta dashes nativamente)
      - cantos arredondados (border-radius: 10px)
    """
    fill = T.C_LT_SLOT if light else T.C_SLOT_BG
    line = T.C_LT_SLOT_LN if light else T.C_SLOT_LINE
    return _add_shape_rect(slide, x, y, w, h, fill_color=fill, line_color=line,
                           line_width=Pt(0.75), rounded=True)


def _pillar_box(slide, x, y, w, h, *, light=False):
    """Caixa estilo .pillar — cantos arredondados, borda accent."""
    fill = T.C_LT_SURF if light else T.C_PILLAR_BG
    line = T.C_ACCENT if light else RGBColor(0x40, 0x10, 0x10)
    return _add_shape_rect(slide, x, y, w, h, fill_color=fill, line_color=line,
                           line_width=Pt(0.5), rounded=True)


def _compare_box(slide, x, y, w, h, *, after=False, light=False):
    """Caixa estilo .compare-col — cantos arredondados."""
    if light:
        fill = RGBColor(0xF0, 0xF0, 0xF0)
        line = RGBColor(0xCC, 0xCC, 0xCC) if not after else T.C_ACCENT
    else:
        fill = T.C_COMPARE_BF if not after else T.C_COMPARE_AF
        line = RGBColor(0x33, 0x33, 0x33) if not after else T.C_ACCENT
    return _add_shape_rect(slide, x, y, w, h, fill_color=fill, line_color=line,
                           line_width=Pt(0.75), rounded=True)


def _ph_icon(slide, label, x, y, size=Inches(0.40), *, light=False):
    """Placeholder .ph-icon: quadrado vermelho com label curto."""
    box = _add_shape_rect(slide, x, y, size, size,
                          fill_color=RGBColor(0x38, 0x0A, 0x0A) if not light else RGBColor(0xFF, 0xEB, 0xEB),
                          line_color=T.C_ACCENT_BR,
                          line_width=Pt(0.5))
    _add_text(slide, label,
              x, y, size, size,
              font_size=Pt(6),
              bold=True,
              color=T.C_ACCENT_BR if not light else T.C_ACCENT_DK,
              align=PP_ALIGN.CENTER)
    return box


def _ph_media(slide, label, x, y, w, h, *, light=False):
    """Placeholder .ph-media: retângulo com rótulo centralizado e cantos arredondados."""
    fill = RGBColor(0x24, 0x24, 0x24) if not light else RGBColor(0xEE, 0xEE, 0xEE)
    line = RGBColor(0x5C, 0x1C, 0x1C)
    _add_shape_rect(slide, x, y, w, h, fill_color=fill, line_color=line,
                    line_width=Pt(0.5), rounded=True)
    _add_text(slide, label,
              x, y, w, h,
              font_size=Pt(8),
              color=T.C_MUTED if not light else T.C_LT_MUTED,
              align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# HTML utilities
# ---------------------------------------------------------------------------

def _has_html(s):
    return bool(re.search(r'<[a-zA-Z]', s or ""))


def _strip_html(s):
    """Remove tags HTML, preserva texto e quebras de linha."""
    s = re.sub(r'<br\s*/?>', '\n', s or "", flags=re.I)
    s = re.sub(r'<li[^>]*>', '\n• ', s, flags=re.I)
    s = re.sub(r'<[^>]+>', '', s)
    s = html.unescape(s)
    s = re.sub(r'\n{3,}', '\n\n', s)
    return s.strip()


def _parse_html_bullets(html_str):
    """
    Extrai itens de <li> de uma string HTML, retornando lista de strings.
    Preserva texto de <strong> com indicador visual.
    """
    items = re.findall(r'<li[^>]*>(.*?)</li>', html_str or "", re.DOTALL | re.I)
    result = []
    for item in items:
        text = _strip_html(item).replace('\n', ' ').strip()
        if text:
            result.append(text)
    return result


def _parse_table_html(html_str):
    """
    Extrai linhas de uma <table> HTML.
    Retorna (headers: list[str], rows: list[list[str]]).
    """
    try:
        clean = re.sub(r'<br\s*/?>', ' ', html_str or "", flags=re.I)
        root = etree.fromstring(f"<root>{clean}</root>")
        headers = []
        rows = []
        for tr in root.iter('tr'):
            cells = []
            for cell in tr:
                if cell.tag in ('th', 'td'):
                    text = ' '.join(
                        (cell.text_content() if hasattr(cell, 'text_content')
                         else etree.tostring(cell, method='text', encoding='unicode'))
                        .split()
                    )
                    cells.append(text)
            if cells:
                if tr.getparent() is not None and tr.getparent().tag == 'thead':
                    headers = cells
                else:
                    rows.append(cells)
        if not headers and rows:
            headers = rows.pop(0)
        return headers, rows
    except Exception:
        return [], []


def _add_text_multirun(slide, runs, x, y, w, h, *,
                       font_name=T.FONT_PRIMARY,
                       default_size=T.PT_BODY,
                       default_color=None,
                       align=PP_ALIGN.LEFT,
                       wrap=True):
    """
    Adiciona caixa de texto com múltiplos runs (suporte bold inline).
    runs: list de dicts { text, bold, color, size }
    """
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.margin_left   = Inches(0.05)
    tf.margin_top    = Inches(0.03)
    tf.margin_right  = Inches(0.05)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    for r in runs:
        run = p.add_run()
        run.text = r.get('text', '')
        f = run.font
        f.name = font_name
        f.size = r.get('size', default_size)
        f.bold = r.get('bold', False)
        c = r.get('color', default_color)
        if c:
            f.color.rgb = c
    return txBox


def _add_bullets_text(slide, items, x, y, w, h, *, light=False):
    """
    Adiciona lista de bullets como caixa de texto.
    items: list[str]
    """
    color = T.C_LT_TEXT if light else T.C_TEXT
    muted = T.C_LT_MUTED if light else T.C_MUTED

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.05)
    tf.margin_top    = Inches(0.03)
    tf.margin_right  = Inches(0.05)
    tf.margin_bottom = Inches(0.03)

    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Bullet character
        r0 = p.add_run()
        r0.text = "• "
        r0.font.name = T.FONT_PRIMARY
        r0.font.size = T.PT_SMALL
        r0.font.color.rgb = T.C_ACCENT_BR
        # Texto
        run = p.add_run()
        run.text = item
        run.font.name = T.FONT_PRIMARY
        run.font.size = T.PT_SMALL
        run.font.color.rgb = color
    return txBox


def _add_table(slide, headers, rows, x, y, w, h, *, light=False):
    """
    Cria tabela nativa python-pptx estilo .deck-proposal-table.
    """
    n_rows = len(rows) + (1 if headers else 0)
    n_cols = max(
        len(headers) if headers else 0,
        max((len(r) for r in rows), default=0)
    )
    if n_rows == 0 or n_cols == 0:
        return

    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table

    # Largura uniforme das colunas
    col_w = w // n_cols
    for ci in range(n_cols):
        tbl.columns[ci].width = col_w

    txt_color = T.C_LT_TEXT if light else T.C_TEXT

    def _cell_style(cell, text, bold=False, fill_color=None, font_size=T.PT_TABLE_BODY):
        cell.text = text
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.name = T.FONT_PRIMARY
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = txt_color
        if fill_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color

    if headers:
        hdr_fill = T.C_LT_TH_BG if light else T.C_TABLE_HEAD
        for ci, h_text in enumerate(headers):
            cell = tbl.cell(0, ci)
            _cell_style(cell, h_text, bold=True, fill_color=hdr_fill,
                        font_size=T.PT_TABLE_HDR)
        row_start = 1
    else:
        row_start = 0

    for ri, row in enumerate(rows):
        fill = (T.C_LT_TR_BG if light else T.C_TABLE_ROW) if ri % 2 == 1 else None
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(row_start + ri, ci)
            _cell_style(cell, cell_text,
                        fill_color=fill or (T.C_LT_BG if light else T.C_BG_DEEP))

    # Bordas finas
    line_c = T.C_LT_TBL_LN if light else T.C_TABLE_LINE
    _set_table_borders(tbl, line_c)


def _set_table_borders(tbl, color: RGBColor):
    """Aplica borda sólida fina a todas as células da tabela."""
    lxml_tbl = tbl._tbl
    hex_c = f"{color.rgb:06X}" if hasattr(color, 'rgb') else f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
    for tr in lxml_tbl.findall(qn('a:tr')):
        for tc in tr.findall(qn('a:tc')):
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn('a:tcPr'))
            for side in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
                ln = tcPr.find(qn(side))
                if ln is None:
                    ln = etree.SubElement(tcPr, qn(side))
                ln.set('w', '6350')  # ~0.5pt em EMU/100
                ln.set('cap', 'flat')
                ln.set('cmpd', 'sng')
                solidFill = ln.find(qn('a:solidFill'))
                if solidFill is None:
                    solidFill = etree.SubElement(ln, qn('a:solidFill'))
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is None:
                    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgbClr.set('val', hex_c)


def _set_grad_bg(slide, light=False):
    """
    Aplica gradiente de fundo ao slide via XML direto.
    Dark: #030303 → #0a0a0a → #141414 em 160deg
    Light: #fafafa → #f4f4f5
    """
    bg = slide.background
    bgPr = bg._element.get_or_add_bgPr()

    if light:
        stops = [("0", "FAFAFA"), ("100000", "F4F4F5")]
        angle = "5400000"   # 90 graus
    else:
        stops = [("0", "030303"), ("50000", "0A0A0A"), ("100000", "141414")]
        angle = "14400000"  # ~160 graus

    # Remove todos os fills existentes antes de adicionar o gradiente
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:pattFill', 'a:blipFill'):
        for el in bgPr.findall(qn(tag)):
            bgPr.remove(el)

    # Constrói gradFill completo
    gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))
    for pos, clr in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', pos)
        srgb = etree.SubElement(gs, qn('a:srgbClr'))
        srgb.set('val', clr)
    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', angle)
    lin.set('scaled', '0')


# ---------------------------------------------------------------------------
# Slide 1 — Abertura (title)
# ---------------------------------------------------------------------------

def render_title(prs, data, logo=None):
    """
    Slide de abertura: accent-bar, título, subtítulo, meta, placeholder hero.
    Espelha: <section class="slide"> com layout-split layout-split--55.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    # Accent bar
    _add_accent_bar(slide, X, Y, w=Inches(0.9), h=Inches(0.055))

    # Slide label
    _add_slide_label(slide, "Abertura", X, Y + Inches(0.10), W)

    # Split 55/45 — col esquerda (55%)
    col_l_w = W * 0.56
    col_r_w = W * 0.40
    col_gap  = W * 0.04
    col_y    = Y + Inches(0.38)
    col_h    = H - Inches(0.38)

    # Título
    title_y = col_y
    title_h = Inches(1.10)
    _add_text(slide,
              data.get('TITLE', ''),
              X, title_y, col_l_w, title_h,
              font_size=T.PT_H1, bold=True,
              color=T.C_TEXT,
              align=PP_ALIGN.LEFT)

    # Subtítulo
    sub_y = title_y + title_h + Inches(0.08)
    sub_h = Inches(0.60)
    _add_text(slide,
              data.get('SUBTITLE', ''),
              X, sub_y, col_l_w, sub_h,
              font_size=T.PT_H2, bold=False,
              color=T.C_MUTED,
              align=PP_ALIGN.LEFT)

    # Meta
    meta_y = sub_y + sub_h + Inches(0.15)
    meta_h = Inches(0.30)
    _add_text(slide,
              data.get('META_LINE', ''),
              X, meta_y, col_l_w, meta_h,
              font_size=Pt(8), color=T.C_MUTED,
              align=PP_ALIGN.LEFT)

    # Placeholder hero (direita)
    ph_x = X + col_l_w + col_gap
    _ph_media(slide, "Placeholder · imagem hero · produto · diagrama",
              ph_x, col_y, col_r_w, col_h - Inches(0.1))

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 2 — Impacto em destaque
# ---------------------------------------------------------------------------

def render_impact(prs, data, logo=None):
    """
    Slide de impacto: big stat em destaque + impact line + ícones KPI.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Impacto em destaque", X, Y, W)

    y_body = Y + Inches(0.28)
    h_body = H - Inches(0.28)

    col_l_w = W * 0.46
    col_r_w = W * 0.49
    col_gap  = W * 0.05

    # Hero stat (esquerda)
    _add_text(slide,
              data.get('BIG_STAT_LABEL', ''),
              X, y_body, col_l_w, Inches(0.35),
              font_size=T.PT_HERO_LABEL, color=T.C_MUTED)

    stat_h = Inches(1.20)
    _add_text(slide,
              data.get('BIG_STAT', ''),
              X, y_body + Inches(0.32), col_l_w, stat_h,
              font_size=T.PT_HERO_STAT, bold=True,
              color=T.C_ACCENT_BR)

    # Direita: impact line + icon row
    r_x = X + col_l_w + col_gap
    r_y = y_body

    # Slot impact line
    slot_h = Inches(1.10)
    _slot_box(slide, r_x, r_y, col_r_w, slot_h)
    _add_text(slide,
              data.get('IMPACT_LINE', ''),
              r_x + Inches(0.08), r_y + Inches(0.06),
              col_r_w - Inches(0.16), slot_h - Inches(0.12),
              font_size=T.PT_BODY, color=T.C_TEXT)

    # Ícones KPI / $ / OP
    icon_y = r_y + slot_h + Inches(0.2)
    icon_size = Inches(0.40)
    icon_gap  = Inches(0.50)
    for i, (lbl, title) in enumerate([("KPI", "Icone KPI"), ("$", "Icone financeiro"), ("OP", "Icone operacao")]):
        ix = r_x + i * (icon_size + icon_gap)
        _ph_icon(slide, lbl, ix, icon_y, icon_size)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 3 — Problema & objetivo (light)
# ---------------------------------------------------------------------------

def render_problem(prs, data, logo=None):
    """
    Slide light — duas colunas: Problema | Objetivo.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=True)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Problema & objetivo", X, Y, W, light=True)

    col_w   = (W - Inches(0.4)) / 2
    col_gap  = Inches(0.4)
    y_body  = Y + Inches(0.30)
    h_body  = H - Inches(0.30)

    # Coluna esquerda — Problema
    cap_h = Inches(0.22)
    _add_text(slide, "Contexto / dor",
              X, y_body, col_w, cap_h,
              font_size=T.PT_MUTED_CAP, color=T.C_LT_MUTED, light=False)

    slot_h = h_body - cap_h - Inches(0.1)
    _slot_box(slide, X, y_body + cap_h + Inches(0.05), col_w, slot_h, light=True)
    _add_text(slide,
              data.get('PROBLEM', ''),
              X + Inches(0.08), y_body + cap_h + Inches(0.1),
              col_w - Inches(0.16), slot_h - Inches(0.14),
              font_size=T.PT_BODY, color=T.C_LT_TEXT)

    # Coluna direita — Objetivo
    rx = X + col_w + col_gap
    _add_text(slide, "Objetivo principal",
              rx, y_body, col_w, cap_h,
              font_size=T.PT_MUTED_CAP, color=T.C_LT_MUTED, light=False)
    _slot_box(slide, rx, y_body + cap_h + Inches(0.05), col_w, slot_h, light=True)
    _add_text(slide,
              data.get('GOAL', ''),
              rx + Inches(0.08), y_body + cap_h + Inches(0.1),
              col_w - Inches(0.16), slot_h - Inches(0.14),
              font_size=T.PT_BODY, color=T.C_LT_TEXT)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 4 — Usuário & decisão Skill vs Interface
# ---------------------------------------------------------------------------

def render_user_decision(prs, data, logo=None):
    """
    Slide dark — duas colunas: Usuário primário | Decisão skill/interface.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Usuário & decisão Skill vs Interface", X, Y, W)

    col_w  = (W - Inches(0.4)) / 2
    col_gap = Inches(0.4)
    y_body = Y + Inches(0.30)
    h_body = H - Inches(0.30)

    icon_size = Inches(0.40)
    icon_gap  = Inches(0.10)
    slot_x_off = icon_size + icon_gap + Inches(0.06)

    # Coluna esquerda
    _ph_icon(slide, "USR", X, y_body, icon_size)
    s_w = col_w - slot_x_off
    slot_h = h_body - Inches(0.1)
    _slot_box(slide, X + slot_x_off, y_body, s_w, slot_h)
    _add_text(slide,
              data.get('USER_PRIMARY', ''),
              X + slot_x_off + Inches(0.08), y_body + Inches(0.08),
              s_w - Inches(0.16), slot_h - Inches(0.16),
              font_size=T.PT_BODY, color=T.C_TEXT)

    # Coluna direita
    rx = X + col_w + col_gap
    _ph_icon(slide, "UX", rx, y_body, icon_size)
    rx2 = rx + slot_x_off
    s_w2 = col_w - slot_x_off

    _add_text(slide, "Decisão (skill / interface / híbrido)",
              rx2, y_body, s_w2, Inches(0.22),
              font_size=T.PT_MUTED_CAP, color=T.C_MUTED)

    cap_h = Inches(0.22)
    slot2_h = h_body - cap_h - Inches(0.12)
    _slot_box(slide, rx2, y_body + cap_h + Inches(0.06), s_w2, slot2_h)
    _add_text(slide,
              data.get('SKILL_OR_UI', ''),
              rx2 + Inches(0.08), y_body + cap_h + Inches(0.12),
              s_w2 - Inches(0.16), slot2_h - Inches(0.16),
              font_size=T.PT_BODY, color=T.C_TEXT)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 5 — Escopo
# ---------------------------------------------------------------------------

def render_scope(prs, data, logo=None):
    """
    Slide dark — Dentro do escopo | Fora do escopo.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Escopo & fora de escopo", X, Y, W)

    col_w  = (W - Inches(0.4)) / 2
    col_gap = Inches(0.4)
    y_body = Y + Inches(0.30)
    h_body = H - Inches(0.30)

    cap_h  = Inches(0.22)
    slot_h = h_body - cap_h - Inches(0.10)

    for i, (key, label) in enumerate([('SCOPE_IN', 'Dentro do escopo'), ('SCOPE_OUT', 'Fora do escopo')]):
        cx = X + i * (col_w + col_gap)
        _add_text(slide, label,
                  cx, y_body, col_w, cap_h,
                  font_size=T.PT_MUTED_CAP, color=T.C_MUTED)
        _slot_box(slide, cx, y_body + cap_h + Inches(0.05), col_w, slot_h)
        _add_text(slide,
                  data.get(key, ''),
                  cx + Inches(0.08), y_body + cap_h + Inches(0.10),
                  col_w - Inches(0.16), slot_h - Inches(0.14),
                  font_size=T.PT_BODY, color=T.C_TEXT)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 6 — Requisitos & critérios (light)
# ---------------------------------------------------------------------------

def render_requirements(prs, data, logo=None):
    """
    Slide light — Bullets de requisitos + slot de métrica de sucesso.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=True)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Requisitos & critérios", X, Y, W, light=True)

    y_cur = Y + Inches(0.32)

    # Caption
    _add_text(slide, "Top requisitos (P0/P1)",
              X, y_cur, W, Inches(0.22),
              font_size=T.PT_MUTED_CAP, color=T.C_LT_MUTED)
    y_cur += Inches(0.24)

    # Bullets
    bullets_raw = data.get('_raw', {}).get('BULLETS_HTML', '')
    items = _parse_html_bullets(bullets_raw)
    if items:
        bullets_h = Inches(1.80)
        _add_bullets_text(slide, items, X, y_cur, W, bullets_h, light=True)
        y_cur += bullets_h + Inches(0.15)

    # Caption métrica
    _add_text(slide, "Critério de sucesso mensurável",
              X, y_cur, W, Inches(0.22),
              font_size=T.PT_MUTED_CAP, color=T.C_LT_MUTED)
    y_cur += Inches(0.24)

    # Slot métrica
    slot_h = min(Inches(1.40), T.SLIDE_H - y_cur - Inches(0.50))
    _slot_box(slide, X, y_cur, W, slot_h, light=True)
    _add_text(slide,
              data.get('SUCCESS_METRIC', ''),
              X + Inches(0.08), y_cur + Inches(0.08),
              W - Inches(0.16), slot_h - Inches(0.16),
              font_size=T.PT_BODY, color=T.C_LT_TEXT)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 7 — Por que agora
# ---------------------------------------------------------------------------

def render_why_now(prs, data, logo=None):
    """
    Slide dark — Why now + timeline.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Por que agora", X, Y, W)

    y_cur = Y + Inches(0.30)

    # Título why now
    title_h = Inches(0.50)
    _add_text(slide,
              data.get('WHY_NOW_TITLE', ''),
              X, y_cur, W, title_h,
              font_size=T.PT_H2, bold=True, color=T.C_TEXT)
    y_cur += title_h + Inches(0.08)

    # Body
    slot_h = Inches(1.10)
    _slot_box(slide, X, y_cur, W, slot_h)
    _add_text(slide,
              data.get('WHY_NOW_BODY', ''),
              X + Inches(0.08), y_cur + Inches(0.08),
              W - Inches(0.16), slot_h - Inches(0.14),
              font_size=T.PT_BODY, color=T.C_TEXT)
    y_cur += slot_h + Inches(0.18)

    # Timeline bullets
    timeline_raw = data.get('_raw', {}).get('TIMELINE_HTML', '')
    items = _parse_html_bullets(timeline_raw)
    if items:
        tl_h = Inches(1.60)
        _add_bullets_text(slide, items, X, y_cur, W, tl_h)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 8 — Solução em pilares
# ---------------------------------------------------------------------------

def render_pillars(prs, data, logo=None):
    """
    Slide dark — 3 pilares em grid horizontal.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Solução em pilares", X, Y, W)

    y_body = Y + Inches(0.32)
    h_body = H - Inches(0.32)
    n_pillars = 3
    gap = Inches(0.20)
    p_w = (W - gap * (n_pillars - 1)) / n_pillars

    for i in range(1, n_pillars + 1):
        px = X + (i - 1) * (p_w + gap)
        _pillar_box(slide, px, y_body, p_w, h_body)

        icon_size = Inches(0.36)
        _ph_icon(slide, f"P{i}", px + Inches(0.08), y_body + Inches(0.10), icon_size)

        # Título do pilar
        title_x = px + icon_size + Inches(0.16)
        title_y = y_body + Inches(0.10)
        title_w = p_w - icon_size - Inches(0.24)
        title_h = Inches(0.40)
        _add_text(slide,
                  data.get(f'PILLAR_{i}_TITLE', '').upper(),
                  title_x, title_y, title_w, title_h,
                  font_size=T.PT_PILLAR_H, bold=True,
                  color=T.C_ACCENT_BR)

        # Descrição
        desc_y = y_body + Inches(0.65)
        desc_h = h_body - Inches(0.75)
        _add_text(slide,
                  data.get(f'PILLAR_{i}_DESC', ''),
                  px + Inches(0.10), desc_y,
                  p_w - Inches(0.20), desc_h,
                  font_size=T.PT_PILLAR_P, color=T.C_MUTED)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 9 — Antes & depois (light)
# ---------------------------------------------------------------------------

def render_before_after(prs, data, logo=None):
    """
    Slide light — compare-col antes | depois.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=True)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Antes & depois", X, Y, W, light=True)

    col_w  = (W - Inches(0.40)) / 2
    col_gap = Inches(0.40)
    y_body = Y + Inches(0.32)
    h_body = H - Inches(0.32)

    icon_size = Inches(0.36)

    for i, (after, title_key, body_key, sym) in enumerate([
        (False, 'BEFORE_TITLE', 'BEFORE_BODY', '−'),
        (True,  'AFTER_TITLE',  'AFTER_BODY',  '+'),
    ]):
        cx = X + i * (col_w + col_gap)
        _compare_box(slide, cx, y_body, col_w, h_body, after=after, light=True)

        # Ícone
        _ph_icon(slide, sym, cx + Inches(0.10), y_body + Inches(0.10), icon_size, light=True)

        # Título
        title_x = cx + icon_size + Inches(0.18)
        title_w = col_w - icon_size - Inches(0.28)
        _add_text(slide,
                  data.get(title_key, '').upper(),
                  title_x, y_body + Inches(0.10), title_w, Inches(0.36),
                  font_size=T.PT_PILLAR_H, bold=True,
                  color=T.C_ACCENT if not after else T.C_ACCENT_BR)

        # Corpo
        body_y = y_body + Inches(0.60)
        body_h = h_body - Inches(0.70)
        _add_text(slide,
                  data.get(body_key, ''),
                  cx + Inches(0.10), body_y, col_w - Inches(0.20), body_h,
                  font_size=Pt(10), color=T.C_LT_MUTED)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 10 — Stack & tecnologia
# ---------------------------------------------------------------------------

def render_tech(prs, data, logo=None):
    """
    Slide dark — tech title + body + badges, placeholder arquitetura.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Stack & tecnologia", X, Y, W)

    col_l_w = W * 0.60
    col_r_w = W * 0.36
    col_gap  = W * 0.04
    y_body  = Y + Inches(0.30)
    h_body  = H - Inches(0.30)

    # Coluna esquerda
    title_h = Inches(0.44)
    _add_text(slide,
              data.get('TECH_TITLE', ''),
              X, y_body, col_l_w, title_h,
              font_size=T.PT_H2, bold=True, color=T.C_TEXT)

    body_y = y_body + title_h + Inches(0.08)
    body_h = Inches(1.20)
    _slot_box(slide, X, body_y, col_l_w, body_h)
    _add_text(slide,
              data.get('TECH_BODY', ''),
              X + Inches(0.08), body_y + Inches(0.08),
              col_l_w - Inches(0.16), body_h - Inches(0.14),
              font_size=T.PT_BODY, color=T.C_TEXT)

    # Badges (parse badges do HTML)
    badges_raw = data.get('_raw', {}).get('TECH_BADGES_HTML', '')
    badge_texts = re.findall(r'>([^<]+)</span>', badges_raw)
    badge_y = body_y + body_h + Inches(0.12)
    bx = X
    for bt in badge_texts:
        badge_w = max(Inches(0.60), Pt(len(bt) * 6))
        badge_h = Inches(0.24)
        _add_shape_rect(slide, bx, badge_y, badge_w, badge_h,
                        fill_color=RGBColor(0x38, 0x10, 0x10),
                        line_color=T.C_ACCENT, line_width=Pt(0.4))
        _add_text(slide, bt, bx + Inches(0.05), badge_y,
                  badge_w - Inches(0.08), badge_h,
                  font_size=Pt(8), color=RGBColor(0xFE, 0xCA, 0xCA))
        bx += badge_w + Inches(0.08)

    # Coluna direita — placeholder
    rx = X + col_l_w + col_gap
    _ph_media(slide, "Arquitetura · print · diagrama C4",
              rx, y_body, col_r_w, h_body * 0.65)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 11 — Métricas & KPIs
# ---------------------------------------------------------------------------

def render_kpi(prs, data, logo=None):
    """
    Slide dark — intro + grid KPI + placeholder dashboard.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Métricas & KPIs", X, Y, W)

    y_cur = Y + Inches(0.30)

    _add_text(slide,
              data.get('KPI_INTRO', ''),
              X, y_cur, W, Inches(0.22),
              font_size=T.PT_MUTED_CAP, color=T.C_MUTED)
    y_cur += Inches(0.25)

    col_l_w = W * 0.58
    col_r_w = W * 0.38
    col_gap  = W * 0.04

    # Grid KPI (parse do HTML: extrai big numbers e labels)
    kpi_raw = data.get('_raw', {}).get('KPI_HTML', '')
    kpi_vals = re.findall(r'<strong[^>]*>([^<]+)</strong>', kpi_raw)
    kpi_labels = re.findall(r'<span[^>]*>([^<]+)</span>', kpi_raw)
    n = len(kpi_vals)

    kpi_h = H - y_cur + T.CONTENT_Y - Inches(0.4)
    card_gap = Inches(0.12)
    n_cols = min(n, 4) if n > 0 else 2
    card_w = (col_l_w - card_gap * (n_cols - 1)) / max(n_cols, 1)
    card_h = kpi_h

    for i, (val, lbl) in enumerate(zip(kpi_vals, kpi_labels)):
        cx = X + i * (card_w + card_gap)
        _add_shape_rect(slide, cx, y_cur, card_w, card_h,
                        fill_color=RGBColor(0x2A, 0x0E, 0x0E),
                        line_color=T.C_ACCENT, line_width=Pt(0.5))
        # Valor grande
        _add_text(slide, val,
                  cx, y_cur + Inches(0.15), card_w, Inches(0.60),
                  font_size=Pt(22), bold=True,
                  color=T.C_ACCENT_BR, align=PP_ALIGN.CENTER)
        # Label
        _add_text(slide, lbl,
                  cx, y_cur + Inches(0.72), card_w, Inches(0.30),
                  font_size=Pt(8), color=T.C_MUTED, align=PP_ALIGN.CENTER)

    # Placeholder dashboard (direita)
    rx = X + col_l_w + col_gap
    _ph_media(slide, "Mini dashboard · print",
              rx, y_cur, col_r_w, kpi_h)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 12 — Stakeholders (light)
# ---------------------------------------------------------------------------

def render_stakeholders(prs, data, logo=None):
    """
    Slide light — título stakeholders + corpo + placeholder foto.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=True)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Pessoas & stakeholders", X, Y, W, light=True)

    col_l_w = W * 0.60
    col_r_w = W * 0.36
    col_gap  = W * 0.04
    y_body  = Y + Inches(0.30)
    h_body  = H - Inches(0.30)

    icon_size = Inches(0.36)
    _ph_icon(slide, "S", X, y_body, icon_size, light=True)

    title_x = X + icon_size + Inches(0.10)
    title_w = col_l_w - icon_size - Inches(0.10)
    _add_text(slide,
              data.get('STAKE_TITLE', ''),
              title_x, y_body, title_w, Inches(0.40),
              font_size=T.PT_H2, bold=True, color=T.C_LT_TEXT)

    slot_y = y_body + Inches(0.48)
    slot_h = h_body - Inches(0.58)
    _slot_box(slide, X, slot_y, col_l_w, slot_h, light=True)
    _add_text(slide,
              data.get('STAKE_BODY', ''),
              X + Inches(0.08), slot_y + Inches(0.08),
              col_l_w - Inches(0.16), slot_h - Inches(0.16),
              font_size=T.PT_BODY, color=T.C_LT_TEXT)

    rx = X + col_l_w + col_gap
    _ph_media(slide, "Foto equipe · logos parceiros",
              rx, y_body, col_r_w, h_body * 0.65, light=True)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 13 — Riscos & próximos passos
# ---------------------------------------------------------------------------

def render_risks(prs, data, logo=None):
    """
    Slide dark — risks slot + placeholder heatmap + next steps slot.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Riscos & próximos passos", X, Y, W)

    y_cur = Y + Inches(0.30)

    col_l_w = W * 0.55
    col_r_w = W * 0.40
    col_gap  = W * 0.05

    # Linha topo: risks + placeholder heatmap
    row_h = Inches(1.30)
    _slot_box(slide, X, y_cur, col_l_w, row_h)
    _add_text(slide,
              data.get('RISKS', ''),
              X + Inches(0.08), y_cur + Inches(0.08),
              col_l_w - Inches(0.16), row_h - Inches(0.14),
              font_size=T.PT_BODY, color=T.C_TEXT)

    rx = X + col_l_w + col_gap
    _ph_media(slide, "Matriz risco · heatmap",
              rx, y_cur, col_r_w, row_h)

    # Próximos passos
    y_cur += row_h + Inches(0.18)
    _add_text(slide, "Próximos passos",
              X, y_cur, W, Inches(0.22),
              font_size=T.PT_MUTED_CAP, color=T.C_MUTED)
    y_cur += Inches(0.24)

    ns_h = H - y_cur + T.CONTENT_Y - Inches(0.40)
    _slot_box(slide, X, y_cur, W, ns_h)
    _add_text(slide,
              data.get('NEXT_STEPS', ''),
              X + Inches(0.08), y_cur + Inches(0.08),
              W - Inches(0.16), ns_h - Inches(0.14),
              font_size=T.PT_BODY, color=T.C_TEXT)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 14 — Fechamento (CTA)
# ---------------------------------------------------------------------------

def render_cta(prs, data, logo=None):
    """
    Slide dark — accent-bar, CTA title, subtitle, body slot.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_accent_bar(slide, X, Y, w=Inches(0.9), h=Inches(0.055))
    _add_slide_label(slide, "Fechamento", X, Y + Inches(0.10), W)

    y_cur = Y + Inches(0.38)

    title_h = Inches(1.10)
    _add_text(slide,
              data.get('CTA_TITLE', ''),
              X, y_cur, W, title_h,
              font_size=T.PT_H1, bold=True, color=T.C_TEXT)
    y_cur += title_h + Inches(0.08)

    sub_h = Inches(0.50)
    _add_text(slide,
              data.get('CTA_SUB', ''),
              X, y_cur, W, sub_h,
              font_size=T.PT_H2, color=T.C_MUTED)
    y_cur += sub_h + Inches(0.18)

    slot_h = H - y_cur + T.CONTENT_Y - Inches(0.40)
    _slot_box(slide, X, y_cur, W, slot_h)
    _add_text(slide,
              data.get('CTA_BODY', ''),
              X + Inches(0.08), y_cur + Inches(0.08),
              W - Inches(0.16), slot_h - Inches(0.14),
              font_size=T.PT_BODY, color=T.C_TEXT)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 15 — Proposta comercial / tabela (light)
# ---------------------------------------------------------------------------

def render_pricing(prs, data, logo=None):
    """
    Slide light — blurb + tabela de preços nativa.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=True)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Proposta comercial — investimento & pacotes",
                     X, Y, W, light=True)

    y_cur = Y + Inches(0.30)
    blurb_h = Inches(0.28)
    _add_text(slide,
              data.get('COMMERCIAL_BLURB', ''),
              X, y_cur, W, blurb_h,
              font_size=T.PT_MUTED_CAP, color=T.C_LT_MUTED)
    y_cur += blurb_h + Inches(0.10)

    # Tabela
    table_html = data.get('_raw', {}).get('PRICING_TABLE_HTML', '')
    headers, rows = _parse_table_html(table_html)
    if headers or rows:
        table_h = H - y_cur + T.CONTENT_Y - Inches(0.35)
        _add_table(slide, headers, rows, X, y_cur, W, table_h, light=True)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Slide 16 — Cenário de receita / chart
# ---------------------------------------------------------------------------

def render_chart(prs, data, logo=None, chart_png=None):
    """
    Slide dark — caption + chart (PNG rasterizado ou placeholder) + legenda.
    chart_png: caminho para PNG do SVG já rasterizado (opcional).
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    _add_slide_label(slide, "Cenário de receita & composição (ilustrativo)", X, Y, W)

    y_cur = Y + Inches(0.30)
    cap_h = Inches(0.22)
    _add_text(slide,
              data.get('CHART_CAPTION', ''),
              X, y_cur, W, cap_h,
              font_size=T.PT_MUTED_CAP, color=T.C_MUTED)
    y_cur += cap_h + Inches(0.08)

    col_l_w = W * 0.60
    col_r_w = W * 0.36
    col_gap  = W * 0.04
    chart_h  = H - y_cur + T.CONTENT_Y - Inches(0.35)

    # Chart (PNG se disponível, senão placeholder)
    if chart_png and Path(chart_png).exists():
        slide.shapes.add_picture(chart_png, X, y_cur, col_l_w, chart_h)
    else:
        _ph_media(slide, "Gráfico SVG (rasterizado)", X, y_cur, col_l_w, chart_h)

    # Legenda
    legend_raw = data.get('_raw', {}).get('CHART_LEGEND_HTML', '')
    legend_items = _parse_legend(legend_raw)
    rx = X + col_l_w + col_gap
    ly = y_cur
    for lbl, clr in legend_items:
        # Swatch
        _add_shape_rect(slide, rx, ly + Inches(0.03), Inches(0.13), Inches(0.13),
                        fill_color=clr or T.C_MUTED)
        _add_text(slide, lbl,
                  rx + Inches(0.18), ly, col_r_w - Inches(0.22), Inches(0.24),
                  font_size=Pt(9), color=T.C_MUTED)
        ly += Inches(0.28)

    _add_logo(slide, logo)
    return slide


def _parse_legend(html_str):
    """Extrai (label, RGBColor) de .deck-chart-legend."""
    items = re.findall(
        r'<span\s+class="swatch"\s+style="background:([^"]+)"[^/]*/>\s*([^<]+)',
        html_str or "", re.I
    )
    result = []
    for clr_str, lbl in items:
        clr_str = clr_str.strip().strip(';')
        try:
            c = T.hex_color(clr_str) if clr_str.startswith('#') else T.C_MUTED
        except Exception:
            c = T.C_MUTED
        result.append((lbl.strip(), c))
    return result


# ---------------------------------------------------------------------------
# Slide 17 — Livre (free HTML) — raster fallback com conteúdo estruturado
# ---------------------------------------------------------------------------

def render_free(prs, data, logo=None, slide_num=1, raster_png=None):
    """
    Slide dark — área livre.
    Se raster_png fornecida: insere imagem full-slide (pixel-perfect).
    Senão: tenta extrair texto/tabela e renderiza o melhor possível.
    """
    slide = _blank_slide(prs)
    _set_grad_bg(slide, light=False)

    X = T.CONTENT_X
    Y = T.CONTENT_Y
    W = T.CONTENT_W
    H = T.CONTENT_H

    label = f"Livre {slide_num} — conteúdo customizado"
    _add_slide_label(slide, label, X, Y, W)

    if raster_png and Path(raster_png).exists():
        slide.shapes.add_picture(raster_png, 0, 0, T.SLIDE_W, T.SLIDE_H)
        _add_logo(slide, logo)
        return slide

    key = f'FREE_HTML_{slide_num}'
    free_html = data.get('_raw', {}).get(key, '')

    y_cur = Y + Inches(0.32)
    body_h = H - Inches(0.42)

    # Tenta tabela primeiro
    headers, rows = _parse_table_html(free_html)
    if headers or rows:
        _add_table(slide, headers, rows, X, y_cur, W, body_h)
        _add_logo(slide, logo)
        return slide

    # Bullets
    items = _parse_html_bullets(free_html)
    if items:
        _add_bullets_text(slide, items, X, y_cur, W, body_h)
        _add_logo(slide, logo)
        return slide

    # Texto puro
    plain = _strip_html(free_html)
    if plain:
        _slot_box(slide, X, y_cur, W, body_h)
        _add_text(slide, plain, X + Inches(0.08), y_cur + Inches(0.08),
                  W - Inches(0.16), body_h - Inches(0.14),
                  font_size=T.PT_BODY, color=T.C_TEXT)

    _add_logo(slide, logo)
    return slide


# ---------------------------------------------------------------------------
# Helper de compatibilidade de _add_text (aceita kwarg light= sem quebrar)
# ---------------------------------------------------------------------------

_orig_add_text = _add_text
def _add_text(slide, text, x, y, w, h, *, light=False, **kw):
    return _orig_add_text(slide, text, x, y, w, h, **kw)
